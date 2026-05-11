# -*- coding: utf-8 -*-
"""슬라이드 헬퍼 — 모든 템플릿이 공유하는 도형/텍스트 빌딩 블록."""
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from .tokens import (
    GREEN, GREEN_PALE, GREEN_BORDER,
    GRAY_900, GRAY_500, GRAY_400, GRAY_100,
    WHITE, FONT_KR,
)

# ── 슬라이드 ──
def new_slide(prs):
    """흰 배경 + 상단 그린 라인이 깔린 빈 슬라이드 추가."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, prs.slide_width, prs.slide_height, fill=WHITE)
    add_rect(slide, 0, 0, prs.slide_width, Inches(0.05), fill=GREEN)
    return slide


# ── 도형 ──
def add_rect(slide, left, top, width, height,
             fill=None, line_color=None, line_width_pt=0.75, corner_pt=None):
    """사각형/둥근사각형 추가. corner_pt 지정 시 라운드(0~0.5 권장)."""
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if corner_pt is not None else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.shadow.inherit = False
    if corner_pt is not None:
        shape.adjustments[0] = corner_pt
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width_pt)
    tf = shape.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.word_wrap = True
    return shape


def add_line_h(slide, left, top, width, color=GREEN, height_emu=20000):
    """수평 가는 선. 그린 언더라인/디바이더용."""
    return add_rect(slide, left, top, width, Emu(height_emu), fill=color)


# ── 텍스트 ──
def set_runs(shape, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.2):
    """단일 paragraph에 여러 run을 세팅. runs 항목: {text, size, bold, color, italic, font, spacing}"""
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    for i, r in enumerate(runs):
        if i == 0 and p.runs:
            run = p.runs[0]
        else:
            run = p.add_run()
        run.text = r["text"]
        run.font.name = r.get("font", FONT_KR)
        run.font.size = Pt(r.get("size", 12))
        run.font.bold = r.get("bold", False)
        run.font.italic = r.get("italic", False)
        if "color" in r:
            run.font.color.rgb = r["color"]
        if r.get("spacing"):
            rPr = run._r.get_or_add_rPr()
            rPr.set("spc", str(r["spacing"]))
    if anchor:
        tf.vertical_anchor = anchor
    return shape


def add_text(slide, left, top, width, height, runs,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.2,
             fill=None, line_color=None, corner_pt=None,
             margin=(0.05, 0.05, 0.05, 0.05)):
    """텍스트 박스(필요시 배경/테두리 포함). margin 단위는 inch."""
    s = add_rect(slide, left, top, width, height,
                 fill=fill, line_color=line_color, corner_pt=corner_pt)
    s.text_frame.margin_left = Inches(margin[0])
    s.text_frame.margin_top = Inches(margin[1])
    s.text_frame.margin_right = Inches(margin[2])
    s.text_frame.margin_bottom = Inches(margin[3])
    set_runs(s, runs, align=align, anchor=anchor, line_spacing=line_spacing)
    return s


def add_paragraphs(shape, paragraphs, margin=(0.2, 0.2, 0.2, 0.2)):
    """여러 paragraph를 한 박스 안에 세팅. paragraphs[i] = list of runs (또는 dict)
    paragraph 단위로 line_spacing/space_before 별도 지정하려면
    paragraphs[i] = {"runs": [...], "line_spacing": 1.4, "space_before": 4}
    """
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin[0])
    tf.margin_top = Inches(margin[1])
    tf.margin_right = Inches(margin[2])
    tf.margin_bottom = Inches(margin[3])

    for idx, item in enumerate(paragraphs):
        if isinstance(item, dict):
            runs = item["runs"]
            ls = item.get("line_spacing", 1.3)
            sb = item.get("space_before")
            align = item.get("align")
        else:
            runs = item
            ls = 1.3
            sb = None
            align = None

        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.line_spacing = ls
        if sb is not None:
            p.space_before = Pt(sb)
        if align is not None:
            p.alignment = align
        for j, r in enumerate(runs):
            if j == 0 and p.runs:
                run = p.runs[0]
            else:
                run = p.add_run()
            run.text = r["text"]
            run.font.name = r.get("font", FONT_KR)
            run.font.size = Pt(r.get("size", 12))
            run.font.bold = r.get("bold", False)
            run.font.italic = r.get("italic", False)
            if "color" in r:
                run.font.color.rgb = r["color"]
            if r.get("spacing"):
                rPr = run._r.get_or_add_rPr()
                rPr.set("spc", str(r["spacing"]))


# ── 칩/배지 ──
def add_pill(slide, left, top, width, height, text,
             *, fill=GREEN_PALE, fg=GREEN, border=None, size=10, bold=True):
    """둥근 칩. fill=GRAY_100, fg=GRAY_400 조합으로 페이지번호용도 가능."""
    s = add_rect(slide, left, top, width, height,
                 fill=fill, line_color=border, corner_pt=0.5)
    set_runs(s, [{"text": text, "size": size, "bold": bold, "color": fg}],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
    return s


# ── 공통 헤더/푸터 ──
def add_header(slide, prs, *, title, eyebrow=None, page=None):
    """모든 슬라이드 공통: 좌측 eyebrow + 메인 타이틀, 우측 페이지번호 칩."""
    if eyebrow:
        add_text(slide, Inches(0.6), Inches(0.45), Inches(8), Inches(0.25),
                 [{"text": eyebrow, "size": 10, "bold": True,
                   "color": GREEN, "spacing": 250}],
                 line_spacing=1.0)
    add_text(slide, Inches(0.6), Inches(0.7), Inches(11), Inches(0.5),
             [{"text": title, "size": 22, "bold": True, "color": GRAY_900}],
             line_spacing=1.0)
    if page:
        add_pill(slide, Inches(12.3), Inches(0.55), Inches(0.95), Inches(0.32),
                 page, fill=GRAY_100, fg=GRAY_400, size=10)


def add_footer(slide, prs,
               left_text="PeopleCore · HR Management System",
               right_text="2팀 최종 프로젝트"):
    y = prs.slide_height - Inches(0.42)
    add_text(slide, Inches(0.6), y, Inches(7), Inches(0.25),
             [{"text": left_text, "size": 9, "color": GRAY_400}],
             anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
    add_text(slide, prs.slide_width - Inches(7) - Inches(0.6), y, Inches(7), Inches(0.25),
             [{"text": right_text, "size": 9, "color": GRAY_400}],
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
